from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import numpy as np

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colour palette
BG = RGBColor(0x0d, 0x11, 0x17)
CARD_BG = RGBColor(0x16, 0x1b, 0x22)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xc9, 0xd1, 0xd9)
MUTED = RGBColor(0x8b, 0x94, 0x9e)
BLUE = RGBColor(0x58, 0xa6, 0xff)
RED = RGBColor(0xf8, 0x51, 0x49)
GREEN = RGBColor(0x56, 0xd3, 0x64)
YELLOW = RGBColor(0xd2, 0x99, 0x22)
PURPLE = RGBColor(0xd2, 0xa8, 0xff)
ORANGE = RGBColor(0xf0, 0x88, 0x3e)

def set_slide_bg(slide, color=BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf

def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_stat_card(slide, left, top, label, value, unit='', accent=BLUE):
    card = add_card(slide, left, top, 2.8, 1.6)
    add_text(slide, left + 0.2, top + 0.15, 2.4, 0.3, label, size=10, color=MUTED, bold=True)
    add_text(slide, left + 0.2, top + 0.5, 2.4, 0.6, f'{value}', size=32, color=WHITE, bold=True)
    add_text(slide, left + 0.2 + len(str(value)) * 0.28, top + 0.65, 1, 0.4, unit, size=14, color=MUTED)

# ════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 1, 1.5, 11, 1.2, 'Warehouse Climate Report', size=44, bold=True, color=WHITE)
add_text(slide, 1, 2.7, 11, 0.8, 'What your sensors are telling you about the office', size=20, color=MUTED)
add_text(slide, 1, 4.0, 11, 0.5, 'Netatmo sensor data  |  NZ Warehouse Office  |  2025', size=14, color=MUTED)
add_text(slide, 1, 5.5, 11, 0.5, 'Based on 17,519 data points across 3 sensor locations', size=12, color=MUTED)

# ════════════════════════════════════════
# SLIDE 2: Key Metrics
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.6, 'The Numbers at a Glance', size=28, bold=True, color=WHITE)

add_stat_card(slide, 0.8, 1.4, 'AVG TEMPERATURE', '20.1', '°C', BLUE)
add_stat_card(slide, 3.8, 1.4, 'AVG HUMIDITY', '70', '%', RED)
add_stat_card(slide, 6.8, 1.4, 'AVG CO2', '496', 'ppm', GREEN)
add_stat_card(slide, 9.8, 1.4, 'CO2 HIGH DAYS', '17', 'days', YELLOW)

# Bottom insight cards
add_card(slide, 0.8, 3.4, 5.6, 3.5)
add_text(slide, 1.0, 3.55, 5.2, 0.4, 'WHAT THE NUMBERS MEAN', size=11, color=MUTED, bold=True)
insights_text = (
    "• Temperature is comfortable year-round (10.9–29.8°C range)\n\n"
    "• Humidity is borderline — 70% avg sits right at mould threshold\n\n"
    "• CO2 averages are fine, but 17 days exceeded 800ppm\n\n"
    "• Worst day hit 1,678 ppm — that's unhealthy air\n\n"
    "• All bad CO2 days happen May–Sept (winter, windows closed)"
)
add_text(slide, 1.0, 3.95, 5.2, 3.0, insights_text, size=13, color=LIGHT)

add_card(slide, 6.8, 3.4, 5.6, 3.5)
add_text(slide, 7.0, 3.55, 5.2, 0.4, 'THE SEASONAL PATTERN', size=11, color=MUTED, bold=True)
seasonal_text = (
    "• Summer (Jan–Feb): Warmest, driest, best air quality\n\n"
    "• Autumn (Mar–May): Cooling down, humidity rising\n\n"
    "• Winter (Jun–Aug): Coldest, wettest, CO2 problems\n\n"
    "• Spring (Sep–Nov): Recovery, most variable weather\n\n"
    "• Dec–Jan: Holiday shutdown, quietest period"
)
add_text(slide, 7.0, 3.95, 5.2, 3.0, seasonal_text, size=13, color=LIGHT)

# ════════════════════════════════════════
# SLIDE 3: The CO2 Problem
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.6, 'The CO2 Problem — When Air Gets Bad', size=28, bold=True, color=WHITE)

# Left: CO2 timeline explanation
add_card(slide, 0.8, 1.3, 6.0, 5.5)
add_text(slide, 1.0, 1.45, 5.6, 0.4, 'DAILY CO2 CYCLE', size=11, color=MUTED, bold=True)
co2_daily = (
    "6:00am  — Lowest point (472 ppm)\n"
    "         Best time to open windows\n\n"
    "7:00am  — People arrive, CO2 starts rising\n\n"
    "1:00pm  — Peak CO2 (514 ppm)\n"
    "         Office is most occupied\n\n"
    "6:00pm  — People leave, CO2 starts falling\n\n"
    "8:00pm  — DANGER ZONE if people still here\n"
    "         Sealed space + no ventilation = CO2 spikes"
)
add_text(slide, 1.0, 1.85, 5.6, 5.0, co2_daily, size=13, color=LIGHT)

# Right: Risk windows
add_card(slide, 7.2, 1.3, 5.4, 2.5)
add_text(slide, 7.4, 1.45, 5.0, 0.4, 'HEALTH RISK EVENTS', size=11, color=RED, bold=True)
risk_text = (
    "11 events exceeded 800ppm\n"
    "54 total hours of poor air quality\n"
    "9 events lasted longer than 3 hours\n\n"
    "Worst: Aug 4–5 — 8.5 hours at 1,678 ppm\n"
    "That's like breathing in a poorly ventilated\n"
    "meeting room all day"
)
add_text(slide, 7.4, 1.85, 5.0, 2.0, risk_text, size=13, color=LIGHT)

# Bottom: Action
add_card(slide, 7.2, 4.1, 5.4, 2.7)
add_text(slide, 7.4, 4.25, 5.0, 0.4, 'WHAT TO DO', size=11, color=GREEN, bold=True)
action_text = (
    "1. Open windows at 6–7am (lowest CO2)\n\n"
    "2. If working past 5pm, open windows again\n\n"
    "3. Consider a CO2 monitor — visual alert when\n"
    "   levels rise above 800ppm\n\n"
    "4. In winter, crack a window even when cold —\n"
    "   fresh air beats stale air for focus"
)
add_text(slide, 7.4, 4.65, 5.0, 2.2, action_text, size=13, color=LIGHT)

# ════════════════════════════════════════
# SLIDE 4: The Humidity Issue
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.6, 'Humidity — The Silent Mould Risk', size=28, bold=True, color=WHITE)

add_card(slide, 0.8, 1.3, 5.6, 2.5)
add_text(slide, 1.0, 1.45, 5.2, 0.4, 'WHERE THE MOISTURE COMES FROM', size=11, color=MUTED, bold=True)
humidity_sources = (
    "43%  Weather infiltration (rain, storms)\n"
    "29%  Early-morning condensation\n"
    "20%  Occupancy (breathing, cooking)\n"
    " 8%  Unknown sources"
)
add_text(slide, 1.0, 1.85, 5.2, 2.0, humidity_sources, size=14, color=LIGHT)

add_card(slide, 6.8, 1.3, 5.6, 2.5)
add_text(slide, 7.0, 1.45, 5.2, 0.4, 'WHEN IT\'S WORST', size=11, color=RED, bold=True)
worst_text = (
    "June: 65% of readings above 75%\n"
    "July: 62% of readings above 75%\n"
    "August: 58% of readings above 75%\n\n"
    "Mould grows above 60% relative humidity\n"
    "Your office averages 70% year-round"
)
add_text(slide, 7.0, 1.85, 5.2, 2.0, worst_text, size=14, color=LIGHT)

add_card(slide, 0.8, 4.2, 11.6, 2.8)
add_text(slide, 1.0, 4.35, 11.2, 0.4, 'RECOMMENDATIONS', size=11, color=GREEN, bold=True)
recs = (
    "• Dehumidifier in Jun–Aug — bring humidity below 60%\n"
    "• Check for water ingress after storms (pressure drops correlate with humidity spikes)\n"
    "• Improve drainage around the building if condensation is morning-only\n"
    "• Consider humidity-resistant materials for any design work stored in the office"
)
add_text(slide, 1.0, 4.75, 11.2, 2.2, recs, size=14, color=LIGHT)

# ════════════════════════════════════════
# SLIDE 5: Building Performance
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.6, 'Building Performance — How the Space Behaves', size=28, bold=True, color=WHITE)

# Three cards
add_card(slide, 0.8, 1.3, 3.8, 2.5)
add_text(slide, 1.0, 1.45, 3.4, 0.4, 'THERMAL MASS', size=11, color=BLUE, bold=True)
add_text(slide, 1.0, 1.85, 3.4, 2.0,
    "Temperature half-life: 9 hours\n\n"
    "Good insulation — heat stays\n"
    "in the building after people\n"
    "leave. Cooling rate: 0.095°C/hr\n\n"
    "Winter heating is 1.8x summer",
    size=13, color=LIGHT)

add_card(slide, 4.9, 1.3, 3.8, 2.5)
add_text(slide, 5.1, 1.45, 3.4, 0.4, 'AIRFLOW', size=11, color=GREEN, bold=True)
add_text(slide, 5.1, 1.85, 3.4, 2.0,
    "No thermal lag between floors\n\n"
    "Air circulates freely between\n"
    "studio, kitchen, and upstairs\n\n"
    "Studio ↔ Kitchen R² = 0.877\n"
    "Studio ↔ Upstairs R² = 0.671",
    size=13, color=LIGHT)

add_card(slide, 9.0, 1.3, 3.8, 2.5)
add_text(slide, 9.2, 1.45, 3.4, 0.4, 'WEATHER SENSITIVITY', size=11, color=YELLOW, bold=True)
add_text(slide, 9.2, 1.85, 3.4, 2.0,
    "11 rapid pressure drops detected\n\n"
    "Storms drop indoor temp 0.66°C\n"
    "and raise humidity 1.7%\n\n"
    "Moderate envelope seal (r=-0.264)",
    size=13, color=LIGHT)

# Bottom energy card
add_card(slide, 0.8, 4.2, 11.6, 2.8)
add_text(slide, 1.0, 4.35, 11.2, 0.4, 'ENERGY PROFILE', size=11, color=ORANGE, bold=True)
energy_text = (
    "Estimated annual heating: ~2,059 kWh    |    Heating degree-days: 858    |    Peak demand: August (6.5 HDD/day)\n\n"
    "The building retains heat well (good insulation), but the kitchen runs 1–2°C warmer than other rooms year-round.\n"
    "This suggests heating is concentrated in the kitchen area — other zones may be underheated in winter."
)
add_text(slide, 1.0, 4.75, 11.2, 2.0, energy_text, size=14, color=LIGHT)

# ════════════════════════════════════════
# SLIDE 6: People Patterns
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.6, 'People Patterns — When the Office is Busy', size=28, bold=True, color=WHITE)

add_card(slide, 0.8, 1.3, 5.6, 2.8)
add_text(slide, 1.0, 1.45, 5.2, 0.4, 'OCCUPANCY PATTERNS', size=11, color=MUTED, bold=True)
occ_text = (
    "Weekday window: 08:00–21:00\n"
    "Peak noise: 3–4pm (54.7 dB)\n"
    "Average occupancy: ~7 people\n\n"
    "Busiest month: February\n"
    "Quietest month: May\n"
    "Weekend activity: 15.1 hrs/day avg"
)
add_text(slide, 1.0, 1.85, 5.2, 2.3, occ_text, size=13, color=LIGHT)

add_card(slide, 6.8, 1.3, 5.6, 2.8)
add_text(slide, 7.0, 1.45, 5.2, 0.4, 'ANOMALIES DETECTED', size=11, color=YELLOW, bold=True)
anomaly_text = (
    "203 anomalous days found\n"
    "14 of 20 worst days were noise-related\n\n"
    "Worst: March 2 — noise z-score 7.5\n"
    "  (something very loud happened)\n\n"
    "Most anomalies are operational,\n"
    "not equipment failures"
)
add_text(slide, 7.0, 1.85, 5.2, 2.3, anomaly_text, size=13, color=LIGHT)

# Bottom: Weekly productivity
add_card(slide, 0.8, 4.5, 5.6, 2.5)
add_text(slide, 1.0, 4.65, 5.2, 0.4, 'WEEKLY PRODUCTIVITY', size=11, color=PURPLE, bold=True)
prod_text = (
    "Winter noisier (47.6 dB) than summer (46.0 dB)\n\n"
    "Busiest: Week 4 (Jan 20) — post-holiday surge\n"
    "Quietest: Week 52 (Dec 22) — Christmas\n\n"
    "Noise follows operational cycles,\n"
    "not seasonal patterns"
)
add_text(slide, 1.0, 5.05, 5.2, 2.0, prod_text, size=13, color=LIGHT)

add_card(slide, 6.8, 4.5, 5.6, 2.5)
add_text(slide, 7.0, 4.65, 5.2, 0.4, 'LEAVE PATTERNS', size=11, color=BLUE, bold=True)
leave_text = (
    "Peak leave: March–April (18 low-activity days)\n"
    "  → Easter + NZ school holidays\n\n"
    "Christmas shutdown: Dec 23–31\n"
    "  → 7 consecutive quiet days\n\n"
    "No leave pattern in Jun–Aug\n"
    "  → People save leave for summer"
)
add_text(slide, 7.0, 5.05, 5.2, 2.0, leave_text, size=13, color=LIGHT)

# ════════════════════════════════════════
# SLIDE 7: Action Items
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.6, 'What You Can Do — Action Items', size=28, bold=True, color=WHITE)

# Priority 1: Health
add_card(slide, 0.8, 1.3, 5.6, 2.5)
add_text(slide, 1.0, 1.45, 5.2, 0.4, '🔴  HEALTH & SAFETY', size=12, color=RED, bold=True)
add_text(slide, 1.0, 1.85, 5.2, 2.0,
    "1. Install a CO2 monitor — visual alert at 800ppm\n"
    "2. Open windows at 6–7am and again after 5pm\n"
    "3. In winter, crack a window when CO2 rises\n"
    "4. Get a dehumidifier for Jun–Aug (mould risk)\n"
    "5. Check for water ingress after storms",
    size=13, color=LIGHT)

# Priority 2: Comfort
add_card(slide, 6.8, 1.3, 5.6, 2.5)
add_text(slide, 7.0, 1.45, 5.2, 0.4, '🟡  COMFORT & PRODUCTIVITY', size=12, color=YELLOW, bold=True)
add_text(slide, 7.0, 1.85, 5.2, 2.0,
    "1. Balance heating across zones (kitchen runs hot)\n"
    "2. Consider sound dampening — noise peaks at 54dB\n"
    "3. Humidity-resistant storage for design materials\n"
    "4. Open-plan ventilation strategy for winter\n"
    "5. Schedule meetings for 6–7am when air is freshest",
    size=13, color=LIGHT)

# Priority 3: Design decisions
add_card(slide, 0.8, 4.2, 11.6, 2.8)
add_text(slide, 1.0, 4.35, 11.2, 0.4, '🟢  FOR YOUR DESIGN WORK', size=12, color=GREEN, bold=True)
design_text = (
    "• Store paper/print materials in low-humidity zones (not near exterior walls in winter)\n"
    "• Consider humidity-resistant finishes for any office upgrades — 70% avg is high\n"
    "• If redesigning the space, prioritise airflow paths between zones (data shows free circulation is good)\n"
    "• CO2 data suggests the office is occupied 08:00–21:00 — design lighting and HVAC for this window\n"
    "• The building retains heat well — passive solar design would amplify this in winter"
)
add_text(slide, 1.0, 4.75, 11.2, 2.2, design_text, size=14, color=LIGHT)

# ════════════════════════════════════════
# SLIDE 8: Closing
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 1, 2.0, 11, 1.0, 'The building is talking.', size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 11, 0.8, 'The sensors are listening.', size=24, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11, 0.8, 'Now you know what to do about it.', size=18, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.0, 11, 0.5, 'Full data analysis available in the interactive dashboard', size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Save
prs.save('/Users/samuel/projects/Netatmo/Warehouse_Climate_Report.pptx')
print('Saved: Warehouse_Climate_Report.pptx')
