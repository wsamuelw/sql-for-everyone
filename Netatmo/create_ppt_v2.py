from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Light theme palette
WHITE = RGBColor(0xff, 0xff, 0xff)
BG_LIGHT = RGBColor(0xf8, 0xf9, 0xfa)
BG_CARD = RGBColor(0xff, 0xff, 0xff)
DARK = RGBColor(0x1a, 0x1a, 0x2e)
TEXT = RGBColor(0x2d, 0x3a, 0x4a)
MUTED = RGBColor(0x6b, 0x7b, 0x8d)
BLUE = RGBColor(0x25, 0x63, 0xeb)
RED = RGBColor(0xdc, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xa3, 0x4a)
YELLOW = RGBColor(0xd9, 0x77, 0x06)
PURPLE = RGBColor(0x7c, 0x3a, 0xed)
ORANGE = RGBColor(0xea, 0x58, 0x0c)
TEAL = RGBColor(0x0d, 0x94, 0x88)

def set_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def text(slide, left, top, width, height, txt, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    return tf

def card(slide, left, top, w, h, fill=WHITE, shadow=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xe5, 0xe7, 0xeb)
    shape.line.width = Pt(1)
    return shape

def accent_bar(slide, left, top, width, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def colour_dot(slide, left, top, color, size=0.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
accent_bar(s, 0, 0, 13.333, BLUE)
text(s, 1.5, 2.2, 10, 1.0, 'What Your Sensors\nAre Telling You', size=42, color=DARK, bold=True, spacing=8)
text(s, 1.5, 3.8, 10, 0.6, 'A year of climate data from the warehouse office', size=18, color=MUTED)
colour_dot(s, 1.5, 5.05, BLUE)
text(s, 1.75, 4.95, 10, 0.5, 'Netatmo sensors  ·  3 locations  ·  17,519 readings  ·  2025', size=12, color=MUTED)

# ════════════════════════════════════════
# SLIDE 2: Headline Insight
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
text(s, 0.8, 0.6, 12, 0.8, 'The Big Picture', size=14, color=MUTED, bold=True)
accent_bar(s, 0.8, 1.0, 2, BLUE)

text(s, 0.8, 1.4, 12, 1.0, 'Your office is comfortable —\nbut the air gets stale in winter.', size=34, color=DARK, bold=True)

# Three stat cards in a row
card(s, 0.8, 3.0, 3.6, 2.0)
accent_bar(s, 0.8, 3.0, 3.6, BLUE)
text(s, 1.1, 3.25, 3.0, 0.3, 'Temperature', size=11, color=MUTED, bold=True)
text(s, 1.1, 3.55, 3.0, 0.7, '20.1°C', size=36, color=DARK, bold=True)
text(s, 1.1, 4.2, 3.0, 0.4, 'Comfortable year-round (10.9–29.8°C)', size=11, color=MUTED)

card(s, 4.8, 3.0, 3.6, 2.0)
accent_bar(s, 4.8, 3.0, 3.6, TEAL)
text(s, 5.1, 3.25, 3.0, 0.3, 'Humidity', size=11, color=MUTED, bold=True)
text(s, 5.1, 3.55, 3.0, 0.7, '70%', size=36, color=DARK, bold=True)
text(s, 5.1, 4.2, 3.0, 0.4, 'Borderline — mould threshold is 60%', size=11, color=RED)

card(s, 8.8, 3.0, 3.6, 2.0)
accent_bar(s, 8.8, 3.0, 3.6, YELLOW)
text(s, 9.1, 3.25, 3.0, 0.3, 'CO2', size=11, color=MUTED, bold=True)
text(s, 9.1, 3.55, 3.0, 0.7, '496 ppm', size=36, color=DARK, bold=True)
text(s, 9.1, 4.2, 3.0, 0.4, 'Fine on average — but 17 bad days', size=11, color=RED)

# Insight callout
card(s, 0.8, 5.5, 11.6, 1.3, fill=RGBColor(0xef, 0xf6, 0xff))
text(s, 1.1, 5.65, 11.0, 0.3, 'KEY INSIGHT', size=10, color=BLUE, bold=True)
text(s, 1.1, 5.95, 11.0, 0.6, 'The building retains heat well, but CO2 builds up when windows are closed. Winter months (May–Sep) account for all 17 bad air quality days.', size=14, color=TEXT)

# ════════════════════════════════════════
# SLIDE 3: CO2 Deep Dive
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
text(s, 0.8, 0.6, 12, 0.8, 'Air Quality', size=14, color=MUTED, bold=True)
accent_bar(s, 0.8, 1.0, 2, RED)

text(s, 0.8, 1.4, 12, 0.8, 'CO2 spikes in winter evenings —\nwhen the building is sealed.', size=32, color=DARK, bold=True)

# Left column - the problem
card(s, 0.8, 2.6, 5.6, 4.2)
text(s, 1.1, 2.8, 5.0, 0.3, 'THE PROBLEM', size=11, color=RED, bold=True)
text(s, 1.1, 3.2, 5.0, 3.5,
    '11 events exceeded 800 ppm\n'
    '54 total hours of poor air quality\n'
    '9 events lasted longer than 3 hours\n\n'
    'Worst event: Aug 4–5\n'
    '8.5 hours at peak 1,678 ppm\n\n'
    'All bad days happen May–Sept\n'
    '(winter, windows closed)',
    size=14, color=TEXT)

# Right column - the pattern
card(s, 6.8, 2.6, 5.6, 4.2)
text(s, 7.1, 2.8, 5.0, 0.3, 'THE DAILY PATTERN', size=11, color=BLUE, bold=True)

# Timeline visual using shapes
hours = [
    ('6:00', '472 ppm', 'Best air — open windows', GREEN),
    ('8:00', '473 ppm', 'People arrive, CO2 starts rising', YELLOW),
    ('13:00', '514 ppm', 'Peak CO2 — most occupied', ORANGE),
    ('18:00', '500 ppm', 'People leave, CO2 falls', YELLOW),
    ('21:00', '508 ppm', 'Risk zone if still occupied', RED),
]
for i, (time, ppm, desc, col) in enumerate(hours):
    y = 3.2 + i * 0.7
    colour_dot(s, 7.2, y + 0.05, col)
    text(s, 7.5, y - 0.05, 1.0, 0.3, time, size=13, color=DARK, bold=True)
    text(s, 8.5, y - 0.05, 1.2, 0.3, ppm, size=13, color=col, bold=True)
    text(s, 9.8, y - 0.05, 2.4, 0.3, desc, size=11, color=MUTED)

# ════════════════════════════════════════
# SLIDE 4: Humidity
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
text(s, 0.8, 0.6, 12, 0.8, 'Humidity', size=14, color=MUTED, bold=True)
accent_bar(s, 0.8, 1.0, 2, TEAL)

text(s, 0.8, 1.4, 12, 0.8, '70% average humidity is a\nyear-round mould risk.', size=32, color=DARK, bold=True)

# Sources breakdown
card(s, 0.8, 2.8, 5.6, 3.8)
text(s, 1.1, 3.0, 5.0, 0.3, 'WHERE THE MOISTURE COMES FROM', size=11, color=MUTED, bold=True)

sources = [
    ('Weather infiltration', '43%', 4.3, BLUE),
    ('Morning condensation', '29%', 2.9, TEAL),
    ('Occupancy', '20%', 2.0, PURPLE),
    ('Other', '8%', 0.8, MUTED),
]
for i, (label, pct, w, col) in enumerate(sources):
    y = 3.5 + i * 0.65
    text(s, 1.1, y, 2.5, 0.3, label, size=12, color=TEXT)
    # Bar
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(y + 0.02), Inches(w * 0.4), Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    text(s, 3.6 + w * 0.4 + 0.15, y, 0.8, 0.3, pct, size=12, color=col, bold=True)

# When it's worst
card(s, 6.8, 2.8, 5.6, 3.8)
text(s, 7.1, 3.0, 5.0, 0.3, 'WHEN IT\'S WORST', size=11, color=MUTED, bold=True)

months_data = [
    ('Jun', 65), ('Jul', 62), ('Aug', 58), ('Sep', 52), ('May', 48)
]
for i, (month, pct) in enumerate(months_data):
    y = 3.5 + i * 0.55
    text(s, 7.1, y, 0.8, 0.3, month, size=13, color=DARK, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(y + 0.02), Inches(pct * 0.05), Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED if pct > 60 else YELLOW
    bar.line.fill.background()
    text(s, 8.0 + pct * 0.05 + 0.15, y, 1.5, 0.3, f'{pct}% of readings >75%', size=11, color=MUTED)

# Callout
card(s, 0.8, 6.0, 11.6, 1.0, fill=RGBColor(0xec, 0xfd, 0xf5))
text(s, 1.1, 6.15, 11.0, 0.3, 'ACTION', size=10, color=GREEN, bold=True)
text(s, 1.1, 6.4, 11.0, 0.4, 'Get a dehumidifier for Jun–Aug. Check for water ingress after storms. Consider humidity-resistant materials for design work.', size=13, color=TEXT)

# ════════════════════════════════════════
# SLIDE 5: Building Performance
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
text(s, 0.8, 0.6, 12, 0.8, 'Building Performance', size=14, color=MUTED, bold=True)
accent_bar(s, 0.8, 1.0, 2, PURPLE)

text(s, 0.8, 1.4, 12, 0.8, 'The building retains heat well —\nbut airflow is the real story.', size=32, color=DARK, bold=True)

# Three insight cards
for i, (title, value, detail, col) in enumerate([
    ('Thermal Mass', '9 hours', 'Temperature half-life after people leave', BLUE),
    ('Airflow', 'No lag', 'Air circulates freely between all 3 floors', GREEN),
    ('Envelope', 'Moderate', 'Storms drop indoor temp 0.66°C over 12hrs', YELLOW),
]):
    left = 0.8 + i * 4.1
    card(s, left, 2.6, 3.8, 2.2)
    accent_bar(s, left, 2.6, 3.8, col)
    text(s, left + 0.3, 2.85, 3.2, 0.3, title, size=11, color=MUTED, bold=True)
    text(s, left + 0.3, 3.2, 3.2, 0.6, value, size=28, color=DARK, bold=True)
    text(s, left + 0.3, 3.8, 3.2, 0.5, detail, size=12, color=MUTED)

# Energy card
card(s, 0.8, 5.2, 11.6, 1.8)
text(s, 1.1, 5.35, 11.0, 0.3, 'ENERGY PROFILE', size=11, color=MUTED, bold=True)
text(s, 1.1, 5.7, 5.0, 0.8,
    'Estimated annual heating: 2,059 kWh\n'
    'Heating degree-days: 858\n'
    'Peak demand: August (6.5 HDD/day)',
    size=13, color=TEXT)
text(s, 6.5, 5.7, 5.5, 0.8,
    'Kitchen runs 1–2°C warmer than other rooms year-round.\n'
    'Heating is concentrated there — other zones may be\n'
    'underheated in winter.',
    size=13, color=TEXT)

# ════════════════════════════════════════
# SLIDE 6: People Patterns
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
text(s, 0.8, 0.6, 12, 0.8, 'People Patterns', size=14, color=MUTED, bold=True)
accent_bar(s, 0.8, 1.0, 2, ORANGE)

text(s, 0.8, 1.4, 12, 0.8, 'The office is active 8am–9pm,\nwith clear seasonal rhythms.', size=32, color=DARK, bold=True)

# Two columns
card(s, 0.8, 2.8, 5.6, 4.0)
text(s, 1.1, 3.0, 5.0, 0.3, 'OCCUPANCY', size=11, color=MUTED, bold=True)

occ_items = [
    ('Active window', '08:00–21:00'),
    ('Peak noise', '3–4pm (54.7 dB)'),
    ('Est. people', '~7 average'),
    ('Busiest month', 'February'),
    ('Quietest month', 'May'),
    ('Weekend activity', '15.1 hrs/day'),
]
for i, (label, value) in enumerate(occ_items):
    y = 3.4 + i * 0.5
    text(s, 1.1, y, 2.2, 0.3, label, size=12, color=MUTED)
    text(s, 3.5, y, 2.5, 0.3, value, size=12, color=DARK, bold=True)

card(s, 6.8, 2.8, 5.6, 4.0)
text(s, 7.1, 3.0, 5.0, 0.3, 'ANOMALIES & LEAVE', size=11, color=MUTED, bold=True)

anomaly_items = [
    ('Anomalous days', '203 detected'),
    ('Worst anomaly', 'Mar 2 (noise z=7.5)'),
    ('Peak leave', 'Mar–Apr (18 days)'),
    ('Christmas', 'Dec 23–31 shutdown'),
    ('Winter leave', 'Almost none'),
    ('Busiest week', 'Jan 20 (post-holiday)'),
]
for i, (label, value) in enumerate(anomaly_items):
    y = 3.4 + i * 0.5
    text(s, 7.1, y, 2.2, 0.3, label, size=12, color=MUTED)
    text(s, 9.5, y, 2.5, 0.3, value, size=12, color=DARK, bold=True)

# ════════════════════════════════════════
# SLIDE 7: Action Items
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
text(s, 0.8, 0.6, 12, 0.8, 'What To Do', size=14, color=MUTED, bold=True)
accent_bar(s, 0.8, 1.0, 2, GREEN)

text(s, 0.8, 1.4, 12, 0.8, 'Three things that would\nmake the biggest difference.', size=32, color=DARK, bold=True)

# Three action cards
actions = [
    ('Ventilate smarter', RED,
     ['Open windows at 6–7am (lowest CO2)',
      'Open again after 5pm if staying late',
      'In winter, crack a window when CO2 rises']),
    ('Control humidity', TEAL,
     ['Get a dehumidifier for Jun–Aug',
      'Check for water ingress after storms',
      'Use humidity-resistant storage for materials']),
    ('Balance the heating', BLUE,
     ['Kitchen runs 1–2°C hotter than other rooms',
      'Consider zoning adjustments',
      'The building retains heat well — use that']),
]

for i, (title, col, items) in enumerate(actions):
    left = 0.8 + i * 4.1
    card(s, left, 2.8, 3.8, 3.8)
    accent_bar(s, left, 2.8, 3.8, col)
    colour_dot(s, left + 0.3, 3.1, col, 0.18)
    text(s, left + 0.6, 3.05, 3.0, 0.3, title, size=16, color=DARK, bold=True)
    for j, item in enumerate(items):
        text(s, left + 0.3, 3.6 + j * 0.6, 3.2, 0.5, f'→  {item}', size=12, color=TEXT)

# ════════════════════════════════════════
# SLIDE 8: Close
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
accent_bar(s, 0, 7.44, 13.333, BLUE)

text(s, 1.5, 2.5, 10, 0.8, 'The building is talking.', size=38, color=DARK, bold=True, align=PP_ALIGN.CENTER)
text(s, 1.5, 3.5, 10, 0.6, 'The sensors are listening.', size=22, color=MUTED, align=PP_ALIGN.CENTER)
text(s, 1.5, 4.8, 10, 0.5, 'Now you know what to do about it.', size=18, color=BLUE, align=PP_ALIGN.CENTER)

text(s, 1.5, 6.0, 10, 0.5, 'Full interactive dashboard and analysis available.', size=12, color=MUTED, align=PP_ALIGN.CENTER)

prs.save('/Users/samuel/projects/Netatmo/Warehouse_Climate_Report.pptx')
print('Saved: Warehouse_Climate_Report.pptx')
